import logging
import json
from litellm import completion
from typing import Optional, Dict, Any, List, Union, Tuple
from fpdf import FPDF
from PIL import Image
import io
import tempfile
import os

logger = logging.getLogger(__name__)

def generate_random_inputs(
    current_title: Optional[str] = None,
    current_description: Optional[str] = None,
    current_style: Optional[str] = None,
    current_pages: Optional[int] = None,
    current_characters: Optional[List[Dict[str, str]]] = None
) -> Optional[Dict[str, Any]]:
    """Generates random story parameters using an LLM, considering existing inputs.

    Args:
        current_title (Optional[str]): The existing title, if provided by the user.
        current_description (Optional[str]): The existing story description, if provided.
        current_style (Optional[str]): The existing image style, if provided.
        current_pages (Optional[int]): The existing number of pages, if provided.
        current_characters (Optional[List[Dict[str, str]]]): The existing characters.

    Returns:
        Optional[Dict[str, Any]]: A dictionary containing the generated story parameters if successful, otherwise None.
    """
    logger.info("Requesting context-aware random inputs generation...")

    context_parts = []
    missing_fields_list = []

    if current_title:
        context_parts.append(f"- Existing Title: {current_title}")
    else:
        missing_fields_list.append("Title")

    if current_description:
        context_parts.append(f"- Existing Story Description: {current_description}")
    else:
        missing_fields_list.append("Story Description")

    if current_style:
        context_parts.append(f"- Existing Image Style: {current_style}")
    else:
        missing_fields_list.append("Image Style")

    context_parts.append(f"- Number of Pages: {current_pages}")

    if current_characters:
        char_context = "\n".join([f"  - {char.get('name', 'Unnamed')}: {char.get('description', 'No description')}" for char in current_characters])
        context_parts.append(f"- Existing Characters:\n{char_context}")
    else:
        missing_fields_list.append("Characters")

    context_str = "\n".join(context_parts)
    missing_str = ", ".join(missing_fields_list) if missing_fields_list else "None (refine existing)"

    prompt = f"""
Here are some existing details provided for a children's storybook:
{context_str}

Based on the context above, please generate creative values *only* for the following missing fields: [{missing_str}].
If no fields are listed as missing, please review the existing details and refine them slightly for better coherence or creativity, while keeping the original spirit.

Follow these guidelines when generating missing information:
1.  Title: Generate a catchy title for the story (string).
2.  Story Description: Generate a short story description (string, 2-4 sentences) that fits the title and characters if they exist.
3.  Image Style: Suggest an appropriate visual style (e.g., Cartoon for children, Watercolor with soft edges, High Quality Digital Painting, Oil Painting with textured brush strokes, Sketch with pencil shading, or Minimalist Line Art). The style should reflect the overall tone of the story described or implied by other fields, such as whimsical, adventurous, or serene.
4.  Characters: If no characters were provided in the context, generate a list of 1 to 3 characters. Each character needs a name (string) and a brief description (string).

Constraints:
- The final "number_of_pages" MUST be exactly {current_pages}.
- If characters were provided in the context, use those exact characters in the output. Do not add or remove characters if they were provided.

Output the *complete* story details, incorporating both the existing information (from the context) and the newly generated information (for the missing fields), as a single JSON object.
The JSON object MUST have keys: "title", "story_description", "image_style", "number_of_pages", "characters" (where characters is a list of objects, each with "name" and "description").
"""

    response_schema = {
        "type": "object",
        "properties": {
            "title": {"type": "string", "description": "Catchy title for the story"},
            "story_description": {"type": "string", "description": "Short story description (2-4 sentences)"},
            "image_style": {"type": "string", "description": "Visual style for images (e.g., Cartoon)"},
            "number_of_pages": {"type": "integer", "description": f"Number of pages (Must be {current_pages})", "minimum": 3, "maximum": 9},
            "characters": {
                "type": "array",
                "description": "List of 1-3 characters (Use existing if provided, else generate)",
                "minItems": 1,
                "maxItems": 3,
                "items": {
                    "type": "object",
                    "properties": {
                        "name": {"type": "string", "description": "Character name"},
                        "description": {"type": "string", "description": "Brief character description"}
                    },
                    "required": ["name", "description"]
                }
            }
        },
        "required": ["title", "story_description", "image_style", "number_of_pages", "characters"]
    }

    logger.info(f"Prompting LLM for fields: {missing_str}")
    # logger.debug(f"Generated Prompt:\n{prompt}") # Optional: Log the full prompt for debugging

    try:
        response = completion(
            model="gemini/gemini-2.0-flash",
            messages=[{"role": "user", "content": prompt}],
            response_format={"type": "json_object", "schema": response_schema},
            temperature=1
        )
        if response and 'choices' in response and response['choices'] and 'message' in response['choices'][0]:
            content = response['choices'][0]['message']['content']
            if isinstance(content, str):
                try:
                    data = json.loads(content)
                    data = data[0] if isinstance(data, list) and len(data) > 0 else data
                    if isinstance(data, dict):
                        # Ensure page number constraint is met
                        if data.get('number_of_pages') != current_pages:
                            logger.warning(f"LLM generated different page number ({data.get('number_of_pages')}) than requested ({current_pages}). Forcing override.")
                            data['number_of_pages'] = current_pages

                        # Ensure character constraint is met (don't modify if user provided chars)
                        if current_characters and data.get('characters') != current_characters:
                             # Check if it *looks* like a reasonable representation, maybe names match?
                             # For now, simpler: if user provided chars, ensure response contains exactly those.
                             # This might be too strict if LLM slightly reformats description.
                             # A more lenient check might be needed, or trust the LLM more.
                             logger.warning(f"LLM modified provided characters. Reverting to original user characters.")
                             data['characters'] = current_characters

                        logger.info(f"Context-aware random data received successfully: {data}")
                        return data
                    else:
                        logger.error(f"Parsed JSON is not a dictionary: {type(data)}. Content: {content}")
                except json.JSONDecodeError as json_err:
                    logger.error(f"Failed to parse JSON response: {json_err}\nRaw content: {content}")
            else:
                logger.error(f"Unexpected content type in response message: {type(content)}")
        else:
            logger.error(f"Unexpected response structure from LLM: {response}")

    except Exception as e:
        logger.exception(f"Error during LLM call for context-aware random inputs: {e}")

    return None # Return None if any error occurs

def sort_image_keys(image_keys: List[str]) -> List[str]:
    """Sorts image keys based on their page numbers.

    Args:
        image_keys: List of image key strings, expected to have page numbers as suffixes.

    Returns:
        List[str]: Sorted list of image keys based on their page numbers.

    Raises:
        ValueError: If keys cannot be properly sorted due to invalid format.
    """
    try:
        return sorted(
            image_keys,
            key=lambda x: int(x.split('_')[-1]) if x.split('_')[-1].isdigit() else 0
        )
    except Exception as e:
        logger.error(f"Error sorting image keys: {e}")
        raise ValueError(f"Failed to sort image keys: {e}")

def process_image_for_pdf(img_data: Union[Image.Image, bytes]) -> Optional[Image.Image]:
    """Processes and converts image data into PDF-compatible format.

    Args:
        img_data: Image data either as PIL Image or bytes.

    Returns:
        Optional[Image.Image]: Processed PIL Image in RGB mode, or None if processing fails.
    """
    try:
        if isinstance(img_data, Image.Image):
            img = img_data
        elif isinstance(img_data, bytes):
            img = Image.open(io.BytesIO(img_data))
        else:
            logger.warning(f"Unsupported image format: {type(img_data)}")
            return None

        # Convert to RGB if needed for PDF compatibility
        if img.mode in ('RGBA', 'P'):
            img = img.convert('RGB')
        return img

    except Exception as e:
        logger.error(f"Failed to process image: {e}")
        return None

def generate_pdf_on_demand(title: str, merged_images_data: Dict[str, Any]) -> Optional[bytes]:
    """Generates the PDF bytes on demand using provided data.

    Args:
        title: The title to use for the PDF.
        merged_images_data: Dictionary of images to include in the PDF.

    Returns:
        Optional[bytes]: The generated PDF bytes, or None if generation fails.
    """
    logger.info("PDF download requested. Generating PDF on demand...")
    
    if not merged_images_data:
        logger.warning("PDF generation attempted but no images found.")
        return None

    images_pil = {}
    try:
        sorted_image_keys = sort_image_keys(list(merged_images_data.keys()))
    except ValueError as e:
        logger.error(f"Error sorting image keys for PDF generation: {e}. Using unsorted keys.")
        sorted_image_keys = list(merged_images_data.keys())

    for key in sorted_image_keys:
        processed_img = process_image_for_pdf(merged_images_data[key])
        if processed_img:
            images_pil[key] = processed_img
        else:
            logger.warning(f"Skipping image for key '{key}' due to processing failure.")

    if not images_pil:
        logger.error("No valid images could be processed for on-demand PDF generation.")
        return None

    try:
        # Use the create_pdf function to generate the PDF
        pdf_bytes = create_pdf(title, images_pil)
        logger.info("On-demand PDF generation successful.")
        return pdf_bytes
    except Exception as e:
        logger.exception(f"Error calling create_pdf during on-demand generation: {e}")
        return None

# --- PDF Generation Utility ---
class PDF(FPDF):
    def header(self):
        # No header needed for storybook pages
        pass

    def footer(self):
        # Page number centered in footer
        self.set_y(-15)
        self.set_font('Arial', 'I', 8)
        self.cell(0, 10, f'Page {self.page_no()}', 0, 0, 'C')

def calculate_optimal_page_dimensions(images_dict: Dict[str, Any]) -> Tuple[float, float]:
    """Calculates optimal page dimensions based on the largest image dimensions.
    
    Args:
        images_dict: Dictionary of images, with keys as page identifiers and values as PIL Images.
        
    Returns:
        Tuple[float, float]: A tuple of (width, height) in mm for optimal page size.
    """
    if not images_dict:
        logger.warning("No images provided for dimension calculation, using default A4")
        return 210, 297  # Default A4 size in mm
    
    # Default minimum dimensions (A4)
    max_width_mm, max_height_mm = 210, 297
    
    # Find the largest image dimensions, excluding cover if possible
    content_images = {k: v for k, v in images_dict.items() if not k.endswith('_0')}
    target_images = content_images if content_images else images_dict
    
    try:
        for img_data in target_images.values():
            if isinstance(img_data, Image.Image):
                # Get dimensions in pixels
                width_px, height_px = img_data.size
                
                # Convert to mm (assuming 96 DPI)
                width_mm = width_px * 25.4 / 96
                height_mm = height_px * 25.4 / 96
                
                # Update max dimensions
                max_width_mm = max(max_width_mm, width_mm)
                max_height_mm = max(max_height_mm, height_mm)
        
        # Add margins (10mm on each side)
        max_width_mm += 20
        max_height_mm += 20
        
        logger.info(f"Calculated optimal page dimensions: {max_width_mm}mm x {max_height_mm}mm")
        return max_width_mm, max_height_mm
    
    except Exception as e:
        logger.exception(f"Error calculating optimal page dimensions: {e}")
        return 210, 297  # Default A4 size in mm

def create_pdf(title: str, images_dict: Dict[str, Any]) -> Optional[bytes]:
    """Creates a PDF document from storybook pages and images.

    Args:
        title: The title of the storybook.
        images_dict: Dictionary mapping page key (e.g., 'page_1')
                                       to image data (expected PIL.Image).

    Returns:
        Optional[bytes]: The generated PDF as bytes if successful, otherwise None.
    """
    logger.info(f"Starting PDF generation for: {title}")
    
    try:
        # Calculate optimal page size
        page_width_mm, page_height_mm = calculate_optimal_page_dimensions(images_dict)
        
        # Create PDF with custom page size
        pdf = PDF(orientation='P', unit='mm', format=(page_width_mm, page_height_mm))
        pdf.set_auto_page_break(auto=True, margin=15)
        pdf.set_left_margin(15)
        pdf.set_right_margin(15)
        pdf.set_title(title)
        pdf.set_author("Storybook Generator")

        # Sort pages based on numeric suffix (e.g., page_1, page_2...)
        try:
            page_keys = sorted(list(images_dict.keys()), key=lambda x: int(x.split('_')[-1]) if x.split('_')[-1].isdigit() else 0)
            if not page_keys:
                logger.error("Cannot generate PDF: No valid page keys found in images_dict.")
                return None
        except Exception as e:
            logger.error(f"Error sorting page keys for PDF generation: {e}")
            return None

        # Create a temporary directory for image files
        temp_dir = tempfile.mkdtemp()
        temp_files = []
        
        try:
            for i, page_key in enumerate(page_keys):
                pdf.add_page()
                is_cover = (i == 0) # Assume first page is cover

                # --- Add Image ---
                img_data = images_dict.get(page_key)
                if img_data:
                    try:
                        # Assuming img_data is a PIL Image object
                        if isinstance(img_data, Image.Image):
                            # Save image to a temporary file
                            temp_img_path = os.path.join(temp_dir, f"temp_img_{i}.png")
                            img_data.save(temp_img_path, format='PNG')
                            temp_files.append(temp_img_path)
                            
                            # Calculate image dimensions and position
                            page_width_mm = pdf.w - pdf.l_margin - pdf.r_margin
                            target_width_ratio = 0.6 if is_cover else 0.95  # Increased ratio for content pages
                            img_width_mm = page_width_mm * target_width_ratio

                            # Get image original aspect ratio
                            orig_w_px, orig_h_px = img_data.size
                            aspect_ratio = orig_h_px / orig_w_px
                            img_height_mm = img_width_mm * aspect_ratio

                            # Center the image horizontally
                            x_pos = pdf.l_margin + (page_width_mm - img_width_mm) / 2
                            y_pos = pdf.get_y() + 10 # Add some top margin

                            # Use the file path with FPDF
                            pdf.image(temp_img_path, x=x_pos, y=y_pos, w=img_width_mm, h=img_height_mm)
                            
                            # Move cursor below the image
                            pdf.set_y(y_pos + img_height_mm + 10)
                        else:
                            logger.warning(f"Skipping image on page {i+1}: Unexpected data type {type(img_data)}")
                            pdf.set_y(pdf.get_y() + 10) # Still add some space

                    except Exception as img_err:
                        logger.error(f"Error processing image for page {i+1} (key: {page_key}): {img_err}")
                        pdf.set_y(pdf.get_y() + 10) # Add space even if image fails
                else:
                    logger.warning(f"No image found for page key: {page_key}")
                    pdf.set_y(pdf.get_y() + 10) # Add space if no image

            # Generate PDF output
            pdf_bytes = pdf.output(dest='S').encode('latin-1') # Output as bytes
            logger.info("PDF generation successful.")
            return pdf_bytes
            
        except Exception as pdf_err:
            logger.exception(f"Error during PDF generation: {pdf_err}")
            return None
            
        finally:
            # Clean up temporary files
            for temp_file in temp_files:
                try:
                    if os.path.exists(temp_file):
                        os.remove(temp_file)
                except Exception as e:
                    logger.warning(f"Failed to remove temporary file {temp_file}: {e}")
                    
            try:
                if os.path.exists(temp_dir):
                    os.rmdir(temp_dir)
            except Exception as e:
                logger.warning(f"Failed to remove temporary directory {temp_dir}: {e}")
    except Exception as e:
        logger.exception(f"Error preparing PDF generation: {e}")
        return None

def create_pptx(title: str, images_dict: Dict[str, Any]) -> Optional[bytes]:
    """Creates a PowerPoint presentation from storybook pages and images.
    
    Args:
        title: The title of the storybook.
        images_dict: Dictionary mapping page key (e.g., 'page_1') to image data (expected PIL.Image).
        
    Returns:
        Optional[bytes]: The generated PPTX as bytes if successful, otherwise None.
    """
    logger.info(f"Starting PPTX generation for: {title}")
    
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        logger.error("python-pptx library not found. Please install with: pip install python-pptx")
        return None
    
    # Sort pages based on numeric suffix
    try:
        page_keys = sort_image_keys(list(images_dict.keys()))
        if not page_keys:
            logger.error("Cannot generate PPTX: No valid page keys found in images_dict.")
            return None
    except Exception as e:
        logger.error(f"Error sorting page keys for PPTX generation: {e}")
        return None
    
    # Calculate optimal page dimensions
    try:
        from pptx.util import Emu
        
        # Calculate optimal slide size based on largest image
        page_width_mm, page_height_mm = calculate_optimal_page_dimensions(images_dict)
        
        # Convert mm to EMU (English Metric Units - used by Office)
        # 1 mm = 36000 EMUs
        width_emu = int(page_width_mm * 36000)
        height_emu = int(page_height_mm * 36000)
        
        # Create new presentation with custom slide size
        prs = Presentation()
        prs.slide_width = width_emu
        prs.slide_height = height_emu
    except Exception as e:
        logger.warning(f"Failed to set custom slide size: {e}. Using default size.")
        prs = Presentation()
    
    # Create a temporary directory for image files
    temp_dir = tempfile.mkdtemp()
    temp_files = []
    
    try:
        # Add slides for each page - using blank slide layout for all slides
        for i, page_key in enumerate(page_keys):
            img_data = images_dict.get(page_key)
            if not img_data or not isinstance(img_data, Image.Image):
                logger.warning(f"Skipping page {page_key}: Invalid image data")
                continue
            
            # Use blank layout for all slides (no placeholders)
            slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Save image to temp file
            temp_img_path = os.path.join(temp_dir, f"temp_img_{i}.png")
            img_data.save(temp_img_path, format='PNG')
            temp_files.append(temp_img_path)
            
            # Calculate image size and position
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            img_width, img_height = img_data.size
            
            # Calculate aspect ratio
            img_ratio = img_height / img_width
            
            # Determine optimal image size based on page number
            is_cover = (i == 0)
            max_width_ratio = 0.8 if is_cover else 0.99  # Cover uses 80% of width, content pages use 95%
            max_width = int(slide_width * max_width_ratio)
            max_height = int(slide_height * 0.99)  # Use 95% of slide height
            
            # Determine constrained dimensions
            if img_ratio > max_height / max_width:  # Height constrained
                height = max_height
                width = height / img_ratio
            else:  # Width constrained
                width = max_width
                height = width * img_ratio
            
            # Calculate center position
            left = (slide_width - width) / 2
            top = (slide_height - height) / 2
            
            # Add the image to the slide (centered)
            slide.shapes.add_picture(temp_img_path, left, top, width, height)
        
        # Save presentation to bytes
        pptx_bytes = io.BytesIO()
        prs.save(pptx_bytes)
        pptx_bytes.seek(0)
        logger.info("PPTX generation successful.")
        return pptx_bytes.getvalue()
    
    except Exception as e:
        logger.exception(f"Error during PPTX generation: {e}")
        return None
    
    finally:
        # Clean up temporary files
        for temp_file in temp_files:
            try:
                if os.path.exists(temp_file):
                    os.remove(temp_file)
            except Exception as e:
                logger.warning(f"Failed to remove temporary file {temp_file}: {e}")
                
        try:
            if os.path.exists(temp_dir):
                os.rmdir(temp_dir)
        except Exception as e:
            logger.warning(f"Failed to remove temporary directory {temp_dir}: {e}")

def generate_pptx_on_demand(title: str, merged_images_data: Dict[str, Any]) -> Optional[bytes]:
    """Generates the PPTX bytes on demand using provided data.
    
    Args:
        title: The title to use for the presentation.
        merged_images_data: Dictionary of images to include in the presentation.
        
    Returns:
        Optional[bytes]: The generated PPTX bytes, or None if generation fails.
    """
    logger.info("PPTX export requested. Generating presentation on demand...")
    
    if not merged_images_data:
        logger.warning("PPTX generation attempted but no images found.")
        return None
    
    try:
        # Process images to ensure they're in the correct format
        images_pil = {}
        sorted_image_keys = sort_image_keys(list(merged_images_data.keys()))
        
        for key in sorted_image_keys:
            processed_img = process_image_for_pdf(merged_images_data[key])  # Reuse the PDF processor
            if processed_img:
                images_pil[key] = processed_img
            else:
                logger.warning(f"Skipping image for key '{key}' due to processing failure.")
        
        if not images_pil:
            logger.error("No valid images could be processed for on-demand PPTX generation.")
            return None
        
        # Use the create_pptx function to generate the presentation
        pptx_bytes = create_pptx(title, images_pil)
        if pptx_bytes:
            logger.info("On-demand PPTX generation successful.")
            return pptx_bytes
        else:
            logger.error("Failed to generate PPTX presentation.")
            return None
    
    except Exception as e:
        logger.exception(f"Error during on-demand PPTX generation: {e}")
        return None